#!/usr/bin/env python3
"""Generate PPT for 1B time reconstruction algorithm."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

BLUE = RGBColor(0x1B, 0x3A, 0x5C)


def title_slide(text):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    tx = s.shapes.add_textbox(Inches(0.4), Inches(0.12), Inches(12.5), Inches(0.55))
    p = tx.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = BLUE
    return s


def txt(s, text, left, top, width, height, size=16):
    tx = s.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tx.text_frame
    tf.word_wrap = True
    for i, line in enumerate(text.split("\n")):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = line
        p.font.size = Pt(size)
        p.space_after = Pt(2)
    return tx


def img(s, path, left, top, width=None, height=None):
    if not os.path.exists(path):
        return
    kw = {}
    if width:
        kw["width"] = Inches(width)
    if height:
        kw["height"] = Inches(height)
    s.shapes.add_picture(path, Inches(left), Inches(top), **kw)


# ═══ 1. 封面 ═══
s = prs.slides.add_slide(prs.slide_layouts[6])
t = txt(s, "HXMT HE Level 1B\n\u65f6\u95f4\u91cd\u5efa\u7b97\u6cd5", 1, 2.2, 11, 2, size=44)
for p in t.text_frame.paragraphs:
    p.alignment = PP_ALIGN.CENTER
    p.font.bold = True
    p.font.color.rgb = BLUE
t2 = txt(s, "\u4e00\u4e2a\u4ece\u8bbe\u8ba1\u4e0a\u5c31\u6297\u9971\u548c\u3001\u6297 CRC \u9519\u8bef\u7684\u65f6\u95f4\u91cd\u5efa\u65b9\u6848", 1, 4.5, 11, 0.8, size=22)
t2.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# ═══ 2. CRC 缺陷与幽灵 ═══
s = title_slide("4-bit CRC \u7684\u7f3a\u9677\uff1a\u5e7d\u7075\u4e8b\u4ef6\u7684\u4ea7\u751f")
img(s, "slides/fig01_crc_ghost.png", 0.3, 0.8, width=12.7)
txt(s, (
    "\u6bcf\u4e2a\u4e8b\u4ef6 8 \u5b57\u8282\uff0c\u6700\u540e 4 bit \u4e3a CRC \u6821\u9a8c\n"
    "\u2022 \u635f\u574f\u6570\u636e\u901a\u8fc7 CRC \u7684\u6982\u7387 = 1/16 = 6.25%\n"
    "\u2022 \u8fd9\u4e9b\u201c\u5e7d\u7075\u201d\u7684 ptime \u548c channel \u662f\u968f\u673a\u503c\uff0c\u4f46\u770b\u8d77\u6765\u548c\u771f\u5b9e\u4e8b\u4ef6\u4e00\u6837\n"
    "\u2022 \u9971\u548c\u671f CRC \u9519\u8bef\u7387\u53ef\u8fbe ~50%\uff0c\u5e7d\u7075\u6570\u91cf\u5de8\u5927"
), 0.4, 5.3, 12, 2, size=16)

# ═══ 3. 核心原则 ═══
s = title_slide("\u6838\u5fc3\u539f\u5219\uff1a\u9519\u8bef\u7684\u662f\u5c11\u6570\uff0cLIS \u627e\u591a\u6570")
img(s, "slides/fig02_principle.png", 0.3, 0.8, width=12.7)
txt(s, (
    "\u4e0d\u80fd\u4fe1\u4efb\u4efb\u4f55\u5355\u4e2a\u4e8b\u4ef6\uff0c\u4f46\u5e7d\u7075\u662f\u5c11\u6570\n"
    "FIFO \u4fdd\u5e8f \u2192 \u771f\u5b9e\u4e8b\u4ef6\u7684 ptime \u4e25\u683c\u5355\u8c03\u9012\u589e\n"
    "\u5e7d\u7075\u7684\u968f\u673a ptime \u6253\u7834\u5355\u8c03\u6027 \u2192 \u88ab\u6700\u957f\u9012\u589e\u5b50\u5e8f\u5217\uff08LIS\uff09\u81ea\u7136\u6392\u9664"
), 0.4, 4.8, 12, 2, size=16)

# ═══ 4. 框架：找 SEC 再填充 ═══
s = title_slide("\u7b97\u6cd5\u6846\u67b6\uff1a\u5148\u627e\u53ef\u4fe1\u7684 SEC \u951a\u70b9\uff0c\u518d\u5728\u951a\u70b9\u4e4b\u95f4\u586b\u5145\u4e8b\u4ef6")
txt(s, (
    "SEC\uff08\u79d2\u4e8b\u4ef6\uff09\uff1a\u786c\u4ef6\u6bcf\u79d2\u6ce8\u5165\u4e00\u4e2a\uff0c\u63d0\u4f9b\u7edd\u5bf9\u65f6\u95f4\u951a\u70b9\n"
    "\u4f46 SEC \u4e5f\u4f1a\u6709\u5e7d\u7075\uff08CRC \u78b0\u649e\u4ea7\u751f\u5047 SEC\uff09\n"
    "\n"
    "\u7b2c\u4e00\u6b65\uff1a\u627e\u5230\u53ef\u4fe1\u7684 SEC\n"
    "  Phase 1 \u2014 \u76f8\u4f4d\u805a\u7c7b\uff1a\u5229\u7528\u786c\u4ef6\u76f8\u4f4d\u5e38\u6570\n"
    "  Phase 2 \u2014 stime LIS\uff1a\u5229\u7528 FIFO \u4fdd\u5e8f\u6027\n"
    "\n"
    "\u7b2c\u4e8c\u6b65\uff1a\u5728\u76f8\u90bb SEC \u4e4b\u95f4\u586b\u5145\u4e8b\u4ef6\n"
    "  \u0394stime=1 \u2192 \u6bcf\u4e8b\u4ef6 1 \u4e2a\u5019\u9009\uff0c\u76f4\u63a5 LIS\n"
    "  \u0394stime>1 \u2192 \u6bcf\u4e8b\u4ef6 k \u4e2a\u5019\u9009\uff0c\u5206\u7ec4 LIS"
), 0.4, 0.8, 12, 6, size=18)

# ═══ 5. 找 SEC Phase 1 ═══
s = title_slide("\u627e SEC Phase 1\uff1a\u76f8\u4f4d\u805a\u7c7b")
img(s, "slides/fig_phase_cluster.png", 0.3, 0.8, width=7.5)
txt(s, (
    "\u771f SEC \u7684\u786c\u4ef6\u76f8\u4f4d\uff1a\n"
    "  phase = (ptime \u2212 stime\u00d7500000)\n"
    "          mod 524288\n"
    "  \u2192 \u96c6\u4e2d\u5728\u7a84\u5e26\uff08\u00b1200 ticks\uff09\n"
    "\n"
    "\u5e7d\u7075 SEC\uff1a\n"
    "  \u968f\u673a stime + \u968f\u673a ptime\n"
    "  \u2192 \u76f8\u4f4d\u5747\u5300\u5206\u5e03\n"
    "\n"
    "\u6ed1\u52a8\u7a97\u53e3\u627e\u6700\u5927\u7c07\n"
    "\u2192 \u81ea\u9002\u5e94\uff0c\u65e0\u9700\u9884\u77e5\u76f8\u4f4d\u503c"
), 8, 0.8, 5, 6, size=16)

# ═══ 6. 找 SEC Phase 2 ═══
s = title_slide("\u627e SEC Phase 2\uff1astime LIS \u5254\u9664\u6f0f\u7f51\u5e7d\u7075")
img(s, "slides/fig_stime_lis.png", 0.3, 0.8, width=7.5)
txt(s, (
    "\u76f8\u4f4d\u805a\u7c7b\u540e\u4ecd\u6709\u5c11\u91cf\u5e7d\u7075\n"
    "\uff08\u78b0\u5de7\u76f8\u4f4d\u4e5f\u5728\u7a84\u5e26\u5185\uff09\n"
    "\n"
    "FIFO \u4fdd\u5e8f\uff1a\n"
    "  \u771f SEC \u7684 stime \u5fc5\u987b\u9012\u589e\n"
    "  \u5e7d\u7075\u7684 stime \u968f\u673a\n"
    "  \u2192 \u6253\u7834\u5347\u5e8f\n"
    "\n"
    "\u5bf9 stime \u5e8f\u5217\u6c42 LIS\n"
    "  \u5e7d\u7075\u88ab\u81ea\u7136\u6392\u9664\n"
    "\n"
    "\u5b9e\u6d4b\uff1a3606 SEC\n"
    "  \u2192 3554 \u6709\u6548 + 52 \u5e7d\u7075"
), 8, 0.8, 5, 6, size=16)

# ═══ 7. ptime 回绕 ═══
s = title_slide("\u586b\u5145\u4e8b\u4ef6\uff1aptime \u56de\u7ed5\u4e0e mod \u89e3\u7b97")
img(s, "slides/fig06_mod_wrap.png", 0.3, 0.8, width=12.7)
txt(s, (
    "MET = SEC.met + elapsed_fwd \u00d7 2\u03bcs + 4.0s\n"
    "elapsed_fwd = (event.ptime \u2212 SEC1.ptime) mod 524288\n"
    "\u2022 mod \u81ea\u52a8\u5904\u7406\u56de\u7ed5\uff0c\u65e0\u9700\u77e5\u9053\u56de\u7ed5\u70b9\u4f4d\u7f6e\n"
    "\u2022 \u0394stime=1 \u65f6\uff0cmod \u7ed3\u679c\u662f\u552f\u4e00\u89e3\uff08PMOD > TICKS_PER_SEC\uff09"
), 0.4, 5.5, 12, 1.5, size=15)

# ═══ 8. 级联死亡 ═══
s = title_slide("\u4e3a\u4ec0\u4e48\u7528 LIS \u800c\u4e0d\u662f\u9010\u4e2a\u6bd4\u8f83\uff1f")
img(s, "slides/fig07_cascade.png", 0.3, 0.8, width=12.7)
txt(s, (
    "\u9010\u4e2a\u6bd4\u8f83\uff1a\u4e00\u4e2a\u5e7d\u7075\u5c31\u80fd\u5236\u9020\u7ea7\u8054\u6b7b\u4ea1\uff08~0.2\u79d2\u6570\u636e\u7a7a\u6d1e\uff09\n"
    "LIS\uff1a\u5168\u5c40\u6700\u4f18\uff0c\u5e7d\u7075\u88ab\u81ea\u7136\u6392\u9664\uff0c\u771f\u5b9e\u4e8b\u4ef6\u4e0d\u53d7\u5f71\u54cd"
), 0.4, 5.5, 12, 1, size=16)

# ═══ 9. Δstime>1 候选 ═══
s = title_slide("\u0394stime>1\uff1a\u6bcf\u4e2a\u4e8b\u4ef6\u6709\u591a\u4e2a\u5019\u9009\u4f4d\u7f6e")
img(s, "slides/fig08_candidates.png", 0.3, 0.8, width=12.7)
txt(s, (
    "\u0394stime=k \u65f6\uff0celapsed = ef + w\u00d7524288\uff0cw = 0,1,...,k-1\n"
    "\u5206\u7ec4 LIS\uff1a\u6bcf\u4e2a\u4e8b\u4ef6\u6700\u591a\u9009\u4e00\u4e2a\u5019\u9009\uff0c\u4f7f\u9009\u51fa\u7684 elapsed \u4e25\u683c\u9012\u589e\uff0c\u6700\u5927\u5316\u9009\u4e2d\u4e8b\u4ef6\u6570"
), 0.4, 5.8, 12, 1, size=16)

# ═══ 10. 降序处理 ═══
s = title_slide("\u5206\u7ec4 LIS \u5173\u952e\uff1a\u540c\u7ec4\u5019\u9009\u964d\u5e8f\u5904\u7406")
img(s, "slides/fig09_descending.png", 0.3, 0.8, width=12.7)
txt(s, (
    "\u5347\u5e8f\uff1a\u5927\u5019\u9009\u5229\u7528\u5c0f\u5019\u9009\u521b\u9020\u7684\u4f4d\u7f6e\uff08\u4f46\u5b83\u4eec\u6765\u81ea\u540c\u4e00\u4e8b\u4ef6\uff01\uff09\n"
    "\u964d\u5e8f\uff1a\u5c0f\u5019\u9009\u53ea\u80fd\u8986\u76d6\u5927\u5019\u9009 \u2192 \u5b89\u5168\uff0c\u4e0d\u4f1a\u81ea\u6211\u5e72\u6270\n"
    "\u590d\u6742\u5ea6\uff1aO(Nk log N)"
), 0.4, 5.5, 12, 1.5, size=16)

# ═══ 11. UTC 约束 ═══
s = title_slide("UTC tail \u7ea6\u675f\uff1a\u7528\u5305\u7684\u6253\u5305\u65f6\u95f4\u526a\u679d")
img(s, "slides/fig_utc_constraint.png", 0.3, 0.8, width=7.5)
txt(s, (
    "\u6bcf\u4e2a CCSDS \u5305\u6709 UTC \u5c3e\n"
    "= MCU \u6253\u5305\u65f6\u7684 GPS \u65f6\u95f4\n"
    "\n"
    "\u4e8b\u4ef6\u4ea7\u751f\u65f6\u95f4 \u2264 UTC tail\n"
    "\u2192 elapsed \u6709\u4e0a\u754c\n"
    "\n"
    "19s \u95f4\u9699\u4e2d\uff1a\n"
    "  \u65e9\u671f\u5305 18\u2192 1 \u4e2a\u5019\u9009\n"
    "  \u540e\u671f\u5305\u7ea6\u675f\u8f83\u677e"
), 8.2, 0.8, 4.8, 6, size=16)

# ═══ 12. 验证 260226A ═══
s = title_slide("\u9a8c\u8bc1\uff1aGRB 260226A \u2014 \u4e0e 1K \u6b8b\u5dee\u4ec5 3 \u4e2a\u4e8b\u4ef6")
img(s, "gap_boundary_grb_260226a_box_a.png", 0.2, 0.8, width=8.5)
txt(s, (
    "GRB 260226A Box A\uff1a\n"
    "  1B: 549,658\n"
    "  1K: 549,661\n"
    "  \u6b8b\u5dee\uff1a3 (0.0005%)\n"
    "\n"
    "\u7f3a\u5931\u539f\u56e0\uff1a\n"
    "  2\u00d7 CRC \u5047\u62d2\u7edd\n"
    "  1\u00d7 dead zone\n"
    "\n"
    "\u84dd = \u0394t=1\n"
    "\u6a59 = \u0394t>1\n"
    "\u7070 = 1K"
), 8.8, 0.8, 4.2, 6, size=16)

# ═══ 13. 验证 221009A ═══
s = title_slide("\u9a8c\u8bc1\uff1aGRB 221009A \u2014 \u6062\u590d 64 \u4e07 SEE \u95f4\u9699\u4e8b\u4ef6")
img(s, "gap_boundary_grb_221009a_box_a.png", 0.2, 0.8, height=2.7)
img(s, "gap_boundary_grb_221009a_box_b.png", 0.2, 3.8, height=2.7)
txt(s, (
    "Box A\uff1a\u5355\u4e2a 19s \u95f4\u9699\n"
    "Box B\uff1a\u591a\u6bb5 2-8s \u95f4\u9699\n"
    "\n"
    "\u4e09 Box \u5149\u53d8\u66f2\u7ebf\u4e00\u81f4\n"
    "\u6062\u590d 64 \u4e07\u4e8b\u4ef6\n"
    "\n"
    "\u6a59\u8272\uff08\u0394t>1\uff09\n"
    "\u586b\u5145 SEE \u95f4\u9699\n"
    "\u4e0e 1K\uff08\u7070\uff09\u57fa\u672c\u543b\u5408"
), 8.5, 0.8, 4.5, 6, size=15)

# ═══ 14. 1K 也会错 ═══
s = title_slide("1K \u4e5f\u4f1a\u72af\u9519\uff1a\u4e0d\u80fd\u8ff7\u4fe1\u6807\u51c6\u7ba1\u7ebf")
img(s, "slides/fig12_1k_wrong.png", 0.3, 0.8, width=8)
txt(s, (
    "GRB 221009A SEE \u95f4\u9699\n"
    "\uff08T+248~270\uff09\n"
    "\n"
    "Box C\uff1a\n"
    "  1K \u4ec5 121,638 \u4e8b\u4ef6\n"
    "  1B \u6709 211,893 \u4e8b\u4ef6\n"
    "  1K \u4e22\u5931 43%\uff01\n"
    "\n"
    "Box A\uff1a\n"
    "  \u603b\u6570\u4e00\u81f4\n"
    "  \u4f46 1K \u8fb9\u754c\u504f\u79fb ~1s"
), 8.5, 0.8, 4.5, 5, size=16)

# ═══ 15. 跨 Box 证明 ═══
s = title_slide("\u8de8 Box \u4e92\u76f8\u5173\uff1a1B \u6bd4 1K \u66f4\u51c6\u786e")
img(s, "slides/fig13_crossbox.png", 0.3, 0.8, width=12.7)
txt(s, (
    "Box B \u6709\u66f4\u5c0f\u7684 SEC \u95f4\u9699\uff08\u0394stime=2~8\uff09\uff0c\u65f6\u95f4\u5b9a\u4f4d\u66f4\u53ef\u9760\uff0c\u4f5c\u4e3a\u53c2\u8003\n"
    "Box A wrap 0\uff08\u5f53\u524d\uff09: |A\u2212B| \u6700\u5c0f \u2192 \u4e0e\u72ec\u7acb Box \u543b\u5408\u6700\u597d\n"
    "Box A wrap 1\uff08+1.05s\uff09: |A\u2212B| \u5dee 3 \u500d \u2192 \u660e\u663e\u66f4\u5dee"
), 0.4, 5.3, 12, 1.5, size=16)

# ═══ 16. 总结 ═══
s = title_slide("\u603b\u7ed3")
txt(s, (
    "\u8bbe\u8ba1\u539f\u5219\uff1a\n"
    "  \u2022 \u4ece\u4e00\u5f00\u59cb\u5c31\u9488\u5bf9\u9971\u548c\u671f\u7684\u6076\u52a3\u6570\u636e\u73af\u5883\u8bbe\u8ba1\n"
    "  \u2022 \u4e0d\u4fe1\u4efb\u4efb\u4f55\u5355\u4e2a\u4e8b\u4ef6\uff0c\u4f46\u9519\u8bef\u7684\u662f\u5c11\u6570\n"
    "  \u2022 FIFO \u4fdd\u5e8f\u6027\uff08LIS\uff09\u662f\u552f\u4e00\u53ef\u9760\u7ea6\u675f\n"
    "\n"
    "\u7b97\u6cd5\u6846\u67b6\uff1a\n"
    "  \u2460 \u627e\u53ef\u4fe1 SEC\uff08\u76f8\u4f4d\u805a\u7c7b + stime LIS\uff09\n"
    "  \u2461 SEC \u95f4\u586b\u5145\u4e8b\u4ef6\uff08\u0394t=1 \u76f4\u63a5 LIS / \u0394t>1 \u5206\u7ec4 LIS + UTC\uff09\n"
    "\n"
    "\u7ed3\u679c\uff1a\n"
    "  \u2022 GRB 260226A\uff1a\u4e0e 1K \u6b8b\u5dee 3 events (0.0005%)\n"
    "  \u2022 GRB 221009A\uff1a\u6062\u590d 64 \u4e07\u4e8b\u4ef6\uff0c\u4e09 Box \u4e00\u81f4\n"
    "  \u2022 \u8de8 Box \u9a8c\u8bc1\uff1a1B \u6bd4 1K \u66f4\u51c6\u786e"
), 0.4, 0.8, 12, 6, size=20)

prs.save("slides/slide_1b_reconstruction.pptx")
print(f"Saved: {len(prs.slides)} slides")
