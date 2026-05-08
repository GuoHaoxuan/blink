"""Shared helpers and color palette for HXMT HE presentation."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Colors (Claude / Anthropic brand palette) ───────────────────────────────
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
DARK = RGBColor(0x2D, 0x28, 0x24)
BLUE = RGBColor(0xD9, 0x77, 0x57)          # Claude terracotta
LIGHT_BLUE = RGBColor(0xF9, 0xEC, 0xE4)
RED = RGBColor(0xBF, 0x3B, 0x30)
LIGHT_RED = RGBColor(0xF6, 0xDD, 0xD6)
GREEN = RGBColor(0x5A, 0x8A, 0x6A)
LIGHT_GREEN = RGBColor(0xE4, 0xEE, 0xE7)
ORANGE = RGBColor(0xC8, 0x7E, 0x33)
LIGHT_ORANGE = RGBColor(0xF8, 0xED, 0xD8)
PURPLE = RGBColor(0x7C, 0x5C, 0x92)
LIGHT_PURPLE = RGBColor(0xEC, 0xE4, 0xF2)
GRAY = RGBColor(0x8C, 0x85, 0x7D)
LIGHT_GRAY = RGBColor(0xF2, 0xEF, 0xEB)
DARK_GRAY = RGBColor(0x5C, 0x55, 0x4E)
BG_COLOR = RGBColor(0xFA, 0xF9, 0xF6)
TITLE_BG = RGBColor(0x2D, 0x28, 0x24)
WARM_ACCENT = RGBColor(0xB0, 0xA8, 0x9E)
WARM_LIGHT = RGBColor(0xD9, 0x9A, 0x7E)
WARM_DESC = RGBColor(0xC8, 0xBA, 0xAD)


def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    return prs


def add_blank_slide(prs):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = BG_COLOR
    return slide


def add_title_bar(prs, slide, title, subtitle=None):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.0))
    bar.fill.solid()
    bar.fill.fore_color.rgb = TITLE_BG
    bar.line.fill.background()
    tf = bar.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.LEFT
    tf.margin_left = Inches(0.5)
    tf.margin_top = Inches(0.15)
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(16)
        p2.font.color.rgb = WARM_ACCENT
        p2.alignment = PP_ALIGN.LEFT


def add_textbox(slide, left, top, width, height, text, font_size=14, bold=False,
                color=None, alignment=PP_ALIGN.LEFT):
    if color is None:
        color = DARK
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment
    return txBox


def add_rich_textbox(slide, left, top, width, height, lines, default_size=14):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if isinstance(line, str):
            text, fs, bold, color = line, default_size, False, DARK
        else:
            text = line[0]
            fs = line[1] if len(line) > 1 else default_size
            bold = line[2] if len(line) > 2 else False
            color = line[3] if len(line) > 3 else DARK
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(fs)
        p.font.bold = bold
        p.font.color.rgb = color
        p.space_after = Pt(4)
    return txBox


def add_box(slide, left, top, width, height, text, fill_color=None,
            border_color=None, font_size=12, bold=False, text_color=None):
    if fill_color is None:
        fill_color = LIGHT_BLUE
    if border_color is None:
        border_color = BLUE
    if text_color is None:
        text_color = DARK
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = border_color
    shape.line.width = Pt(1.5)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(6)
    tf.margin_right = Pt(6)
    tf.margin_top = Pt(4)
    tf.margin_bottom = Pt(4)
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = text_color
    p.alignment = PP_ALIGN.CENTER
    return shape


def _line_shape(slide, x1, y1, x2, y2, color, width):
    """Create a line using a thin rectangle (avoids connector XML issues)."""
    from pptx.util import Emu
    dx = int(x2 - x1)
    dy = int(y2 - y1)
    adx = abs(dx)
    ady = abs(dy)
    min_dim = int(Emu(2000))  # minimum visible dimension

    if ady < min_dim and adx >= min_dim:
        # Horizontal line
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                        int(min(x1, x2)), int(y1), adx, int(width))
    elif adx < min_dim and ady >= min_dim:
        # Vertical line
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                        int(x1), int(min(y1, y2)), int(width), ady)
    else:
        # Diagonal — use a thin rectangle at an angle (approximate)
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                        int(min(x1, x2)), int(min(y1, y2)),
                                        max(adx, min_dim), max(ady, min_dim))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_arrow(slide, x1, y1, x2, y2, color=None, width=Pt(2)):
    """Draw a line with an arrowhead triangle at (x2, y2)."""
    if color is None:
        color = DARK_GRAY
    import math
    # Draw the line body
    shape = _line_shape(slide, int(x1), int(y1), int(x2), int(y2), color, width)
    # Draw arrowhead as a small triangle at the end
    sz = int(Inches(0.12))
    dx = int(x2) - int(x1)
    dy = int(y2) - int(y1)
    length = math.sqrt(dx * dx + dy * dy)
    if length > 0:
        tri = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE,
                                      int(x2) - sz // 2,
                                      int(y2) - sz // 2,
                                      sz, sz)
        tri.fill.solid()
        tri.fill.fore_color.rgb = color
        tri.line.fill.background()
        angle_deg = math.degrees(math.atan2(dx, -dy))
        tri.rotation = angle_deg
    return shape


def add_line(slide, x1, y1, x2, y2, color=None, width=Pt(1), dashed=False):
    if color is None:
        color = GRAY
    shape = _line_shape(slide, int(x1), int(y1), int(x2), int(y2), color, width)
    # dashed is ignored for rectangle-based lines (visual approximation)
    return shape


def draw_bar(slide, left, top, width, height, fill_color, label="", font_size=10,
             text_color=None, border_color=None):
    if text_color is None:
        text_color = WHITE
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = fill_color
    if border_color:
        bar.line.color.rgb = border_color
        bar.line.width = Pt(0.5)
    else:
        bar.line.fill.background()
    if label:
        tf = bar.text_frame
        tf.word_wrap = False
        tf.margin_left = Pt(3)
        tf.margin_right = Pt(3)
        tf.margin_top = Pt(1)
        tf.margin_bottom = Pt(1)
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(font_size)
        p.font.color.rgb = text_color
        p.alignment = PP_ALIGN.CENTER
    return bar


def draw_table(slide, left, top, col_widths, rows, row_height=Inches(0.4),
               header_bg=None, header_text_color=None, font_size=12):
    if header_bg is None:
        header_bg = BLUE
    if header_text_color is None:
        header_text_color = WHITE
    for r, row in enumerate(rows):
        y = top + r * row_height
        x = left
        for c, (text, cw) in enumerate(zip(row, col_widths)):
            cell = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, cw, row_height)
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_bg
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if r % 2 == 1 else LIGHT_GRAY
            cell.line.color.rgb = GRAY
            cell.line.width = Pt(0.5)
            tf = cell.text_frame
            tf.margin_left = Pt(4)
            tf.margin_right = Pt(4)
            tf.margin_top = Pt(2)
            tf.margin_bottom = Pt(2)
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = text
            p.font.size = Pt(font_size)
            p.font.color.rgb = header_text_color if r == 0 else DARK
            p.font.bold = r == 0
            p.alignment = PP_ALIGN.CENTER
            x += cw


def draw_sawtooth(slide, origin_x, origin_y, plot_w, plot_h, n_cycles=3,
                  line_color=None, axis_color=None):
    if line_color is None:
        line_color = BLUE
    if axis_color is None:
        axis_color = GRAY
    add_line(slide, origin_x, origin_y, origin_x, origin_y - plot_h,
             color=axis_color, width=Pt(1.5))
    add_line(slide, origin_x, origin_y, origin_x + plot_w, origin_y,
             color=axis_color, width=Pt(1.5))
    add_arrow(slide, origin_x, origin_y - plot_h + Inches(0.2),
              origin_x, origin_y - plot_h, color=axis_color, width=Pt(1.5))
    add_arrow(slide, origin_x + plot_w - Inches(0.2), origin_y,
              origin_x + plot_w, origin_y, color=axis_color, width=Pt(1.5))
    cycle_w = int((plot_w - Inches(0.3)) / n_cycles)
    for i in range(n_cycles):
        x_start = origin_x + cycle_w * i
        x_end = origin_x + cycle_w * (i + 1)
        y_bottom = origin_y
        y_top = origin_y - plot_h + Inches(0.15)
        add_line(slide, x_start, y_bottom, x_end, y_top,
                 color=line_color, width=Pt(2.5))
        if i < n_cycles - 1:
            add_line(slide, x_end, y_top, x_end, y_bottom,
                     color=line_color, width=Pt(2.5))
    return cycle_w
